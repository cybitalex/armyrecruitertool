#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for Army Recruiter Tool
Simple, concise presentation for Colonel review - leads to product demo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Army color palette (RGB values)
ARMY_BLACK = RGBColor(34, 31, 32)
ARMY_GOLD = RGBColor(255, 204, 1)
ARMY_GREEN = RGBColor(47, 55, 47)
ARMY_TAN = RGBColor(241, 228, 199)
ARMY_FIELD01 = RGBColor(114, 115, 101)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# Logo paths
ARMY_LOGO_PATH = "army_logo1_rev_rgb_300ppi.png"  # For header/footer (dark backgrounds)
ARMY_LOGO_TITLE_PATH = "army_logo1_pos_rgb_300ppi.png"  # For title slide (light background)
AIRBORNE_LOGO_PATH = "airborne-inn-logo-white 2.png"

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_slide(slide)
    
    # Slide 2: The Challenge
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout for custom header/footer
    add_content_slide_with_header_footer(slide, "The Challenge", [
        "Manual tracking and administrative burden",
        "Difficulty identifying prime recruitment locations",
        "Limited visibility into performance metrics",
        "Inefficient prospect capture methods"
    ])
    
    # Slide 3: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_with_header_footer(slide, "The Solution", [
        "All-in-one web-based recruitment platform",
        "Real-time dashboard with key metrics",
        "Station commander oversight and reporting",
        "QR code system for instant prospect capture",
        "AI-powered prospecting map for location intelligence",
        "Event management and automated communications",
        "",
        "Ready for immediate deployment"
    ])
    
    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_with_header_footer(slide, "Key Features", [
        "Dashboard: Real-time recruitment statistics and tracking",
        "Station Commander View: Aggregated team metrics and exports",
        "QR Codes: Instant prospect capture and attribution",
        "Prospecting Map: AI-powered location discovery with scores",
        "Event Management: Track and plan recruitment events",
        "Automated Emails: Professional, Army-themed communications",
        "Survey System: Real-time presentation feedback"
    ])
    
    # Slide 5: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_with_header_footer(slide, "Operational Benefits", [
        "Reduces administrative time by 50%",
        "Increases recruiter efficiency and effectiveness",
        "Data-driven decision making",
        "Better targeting of prime locations",
        "Professional, modern image",
        "",
        "More recruits, less paperwork"
    ])
    
    # Slide 6: Security & Status
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_with_header_footer(slide, "Security & Deployment Status", [
        "Enterprise-grade security:",
        "  • HTTPS/SSL encryption",
        "  • Secure authentication and password hashing",
        "  • Role-based access control (Recruiter/Commander/Admin)",
        "  • Protected routes and input validation",
        "",
        "Production Ready:",
        "  • Fully functional and deployed",
        "  • Scalable architecture"
    ])
    
    # Slide 7: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_with_header_footer(slide, "Next Steps", [
        "Approve pilot program for select stations",
        "Schedule training for users",
        "Establish success metrics",
        "Begin phased rollout",
        "",
        "Ready to proceed upon approval"
    ])
    
    # Slide 8: Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_demo_slide(slide)
    
    return prs

def add_header_footer(slide):
    """Add Army-themed header and footer to a slide"""
    # Header bar
    header_left = Inches(0)
    header_top = Inches(0)
    header_width = Inches(10)
    header_height = Inches(0.8)
    
    header_shape = slide.shapes.add_shape(1, header_left, header_top, header_width, header_height)  # Rectangle
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = ARMY_BLACK
    header_shape.line.fill.background()
    
    # Army logo in header (left)
    if os.path.exists(ARMY_LOGO_PATH):
        try:
            army_logo = slide.shapes.add_picture(ARMY_LOGO_PATH, Inches(0.2), Inches(0.1), height=Inches(0.6))
        except:
            pass
    
    # Airborne logo in header (right)
    if os.path.exists(AIRBORNE_LOGO_PATH):
        try:
            airborne_logo = slide.shapes.add_picture(AIRBORNE_LOGO_PATH, Inches(8.5), Inches(0.1), height=Inches(0.6))
        except:
            pass
    
    # Footer bar
    footer_left = Inches(0)
    footer_top = Inches(6.7)
    footer_width = Inches(10)
    footer_height = Inches(0.8)
    
    footer_shape = slide.shapes.add_shape(1, footer_left, footer_top, footer_width, footer_height)  # Rectangle
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = ARMY_GREEN
    footer_shape.line.fill.background()
    
    # Footer text
    footer_text_left = Inches(0.2)
    footer_text_top = Inches(6.85)
    footer_text_width = Inches(9.6)
    footer_text_height = Inches(0.5)
    footer_text_box = slide.shapes.add_textbox(footer_text_left, footer_text_top, footer_text_width, footer_text_height)
    footer_frame = footer_text_box.text_frame
    footer_frame.text = "Army Recruiter Tool | Proprietary Software | Copyright © 2025 Alex Moran"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = ARMY_TAN
    footer_para.alignment = PP_ALIGN.CENTER

def add_title_slide(slide):
    """Add title slide with logos"""
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header/footer
    add_header_footer(slide)
    
    # Army logo (left side, larger) - use positive version for light background
    if os.path.exists(ARMY_LOGO_TITLE_PATH):
        try:
            army_logo = slide.shapes.add_picture(ARMY_LOGO_TITLE_PATH, Inches(1), Inches(2), height=Inches(1.5))
        except:
            pass
    
    # Airborne logo (right side, larger)
    if os.path.exists(AIRBORNE_LOGO_PATH):
        try:
            airborne_logo = slide.shapes.add_picture(AIRBORNE_LOGO_PATH, Inches(7), Inches(2), height=Inches(1.5))
        except:
            pass
    
    # Title
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Army Recruiter Tool"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    top = Inches(5.2)
    height = Inches(0.6)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Modernizing Recruitment Operations"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = BLACK
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide_with_header_footer(slide, title_text, bullet_points):
    """Add a content slide with header, footer, title and bullet points"""
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header and footer
    add_header_footer(slide)
    
    # Title
    title_left = Inches(0.5)
    title_top = Inches(1.2)
    title_width = Inches(9)
    title_height = Inches(0.8)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.LEFT
    
    # Content area with bullet points
    content_left = Inches(0.8)
    content_top = Inches(2.2)
    content_width = Inches(8.4)
    content_height = Inches(4.2)
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        if point.strip():  # Only add bullet if not empty
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.space_after = Pt(10)
        else:
            # Empty line for spacing
            p.text = " "
            p.space_after = Pt(5)

def add_demo_slide(slide):
    """Slide 8: Demo"""
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header/footer
    add_header_footer(slide)
    
    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Live Demo"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.CENTER
    
    # URL
    top = Inches(4.2)
    height = Inches(0.8)
    url_box = slide.shapes.add_textbox(left, top, width, height)
    url_frame = url_box.text_frame
    url_frame.text = "https://armyrecruitertool.duckdns.org"
    url_para = url_frame.paragraphs[0]
    url_para.font.size = Pt(24)
    url_para.font.color.rgb = BLACK
    url_para.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    print("Generating simplified PowerPoint presentation with logos and headers...")
    prs = create_presentation()
    output_file = "Army_Recruiter_Tool_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")
    print("Ready for demo!")
